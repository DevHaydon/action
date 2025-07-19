from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

import os
from datetime import datetime

# ----------------------------------------------------------------------------------
# CONFIG
# ----------------------------------------------------------------------------------
SLIDE_W = Emu(13.333 * 914400/1)  # Actually we'll override using slidesize
SLIDE_H = Emu(7.5 * 914400/1)

# Brand colours (RGB)
GROUND_FOREST = RGBColor(0x1B, 0x5E, 0x20)
CANOPY_LIGHT = RGBColor(0x66, 0xBB, 0x6A)
VOLT_FIELD = RGBColor(0xCE, 0xFF, 0x00)
HIVIS_AMBER = RGBColor(0xFF, 0x6F, 0x00)
SOIL_BROWN = RGBColor(0x5D, 0x40, 0x37)
CARBON_BLACK = RGBColor(0x00, 0x00, 0x00)
CHALK_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OVERCAST = RGBColor(0xF5, 0xF5, 0xF5)
TRACKLINE_GREY = RGBColor(0x9E, 0x9E, 0x9E)
PERFORMANCE_RED = RGBColor(0xD5, 0x00, 0x00)

# Margins (inches)
MARGIN_L = Inches(0.67)  # ~48px at 72dpi approx
MARGIN_T = Inches(0.67)
MARGIN_R = Inches(0.67)
MARGIN_B = Inches(0.67)

# Fonts (PowerPoint stores font names; it will substitute if not installed)
FONT_HEAD = "Barlow Condensed ExtraBold"
FONT_SUB = "Montserrat SemiBold"
FONT_BODY = "Inter"
FONT_FALLBACK_HEAD = "Arial Narrow"
FONT_FALLBACK_BODY = "Arial"

# ----------------------------------------------------------------------------------
# Utilities
# ----------------------------------------------------------------------------------
def set_shape_fill(shape, rgb):
    """Set solid fill colour on a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def set_shape_line(shape, rgb=None, width_pt=0):
    line = shape.line
    if width_pt <= 0:
        line.fill.background()
        line.width = Pt(0)
    else:
        line.width = Pt(width_pt)
        line.fill.solid()
        line.fill.fore_color.rgb = rgb or CARBON_BLACK

def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY, 
                font_size=Pt(20), bold=False, colour=CARBON_BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = font_size
    font.bold = bold
    font.color.rgb = colour
    p.alignment = align
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    return txBox

def add_title(slide, text, colour=CARBON_BLACK):
    """Standard content slide H3 title style with Volt underline rule."""
    left = MARGIN_L
    top = MARGIN_T
    width = slide.part.slide_layout.slide_master.part.width - (MARGIN_L + MARGIN_R)
    height = Pt(60)
    box = add_textbox(slide, left, top, width, height, text, font_name=FONT_SUB, font_size=Pt(40), bold=True, colour=colour)
    # underline rule
    # width matches text box width; thin rectangle
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top + Pt(46), width, Pt(4))
    set_shape_fill(shape, VOLT_FIELD)
    set_shape_line(shape, None, 0)
    return box

def add_logo_lockup(slide, left=None, top=None, height=Pt(24)):
    """Create placeholder lockup: Ground Control | Nike"""
    if left is None:
            left = prs.slide_width - MARGIN_R - Inches(3)
    if top is None:
            top = prs.slide_height - MARGIN_B + Pt(4)  # slight raise above footer maybe
    width = Inches(3)

    # container group? python-pptx lacks grouping until after shapes exist; we'll just sequentially draw
    # GC placeholder
    gc_box = slide.shapes.add_textbox(left, top, width/2.2, height)
    tf = gc_box.text_frame
    tf.text = "GROUND CONTROL LOGO"
    tf.paragraphs[0].runs[0].font.size = Pt(10)
    tf.paragraphs[0].runs[0].font.name = FONT_BODY
    tf.paragraphs[0].runs[0].font.color.rgb = TRACKLINE_GREY
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # divider
    div_left = left + width/2.2 + Pt(4)
    div = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, div_left, top + Pt(6), Pt(1), height - Pt(12))
    set_shape_fill(div, TRACKLINE_GREY)
    set_shape_line(div, None, 0)

    # Nike placeholder
    nike_left = div_left + Pt(4)
    nike_box = slide.shapes.add_textbox(nike_left, top, width/2.4, height)
    tf2 = nike_box.text_frame
    tf2.text = "NIKE SWOOSH"
    tf2.paragraphs[0].runs[0].font.size = Pt(10)
    tf2.paragraphs[0].runs[0].font.name = FONT_BODY
    tf2.paragraphs[0].runs[0].font.color.rgb = TRACKLINE_GREY
    tf2.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_footer(slide):
    # simple page number & lockup
    add_logo_lockup(slide)
    # page number placeholder
    pn_box = slide.shapes.add_textbox(MARGIN_L, slide.height - MARGIN_B + Pt(4), Inches(0.8), Pt(24), "<<#>>",
                                      font_name=FONT_BODY, font_size=Pt(10), bold=False, colour=TRACKLINE_GREY)
    # Actually can't auto field easily; so show placeholder text; user can insert slide number in master later.

def add_groundline_swoosh(slide, colour_start=GROUND_FOREST, colour_end=CANOPY_LIGHT, opacity=1.0, 
                          rise_fraction=0.65):
    """
    Add a rising angled band from bottom-left to near top-right at ~23°.
    We'll draw a freeform polygon wider than slide so rotation cropping works.
    """
    w = slide.width
    h = slide.height
    # coordinates: bottom-left outside, bottom-left margin, top-right fraction, top-right outside
    # Use freeform builder
    fb = slide.shapes.build_freeform(colour_start, 0, 0)  # colour_start is not used in build
    # Actually build_freeform expects starting coordinates; We'll approximate polygon points manually via xml? easier: use parallelogram rotated rectangle.

    # Simpler: draw rectangle across full slide, rotate.
    band_height = Emu(h * 0.25)  # 25% height band
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, h - band_height, w, band_height)
    shape.rotation = -23  # negative to rise L->R
    fill = shape.fill
    fill.gradient()
    grad = fill.gradient_stops
    grad[0].color.rgb = colour_start
    grad[1].color.rgb = colour_end
    # gradient to Volt tip small stop
    gs3 = grad.add_stop_position(0.85)
    gs3.color.rgb = VOLT_FIELD
    set_shape_line(shape, None, 0)
    return shape

def add_scrim(slide, rgb=GROUND_FOREST, opacity=0.7):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill = rect.fill
    fill.solid()
    fill.fore_color.rgb = rgb
    fill.fore_color.transparency = 1 - opacity
    set_shape_line(rect, None, 0)
    rect.z_order = 1
    return rect

def add_placeholder_image_frame(slide, left, top, width, height, label="Image"):
    ph = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    set_shape_fill(ph, TRACKLINE_GREY)
    set_shape_line(ph, CHALK_WHITE, 2)
    tx = ph.text_frame
    tx.text = label
    tx.paragraphs[0].alignment = PP_ALIGN.CENTER
    tx.paragraphs[0].runs[0].font.color.rgb = CHALK_WHITE
    tx.paragraphs[0].runs[0].font.bold = True
    return ph

# ----------------------------------------------------------------------------------
# PRESENTATION INIT
# ----------------------------------------------------------------------------------
prs = Presentation()
# set slide size 16:9 1920x1080
prs.slide_width = Emu(1920 * 9525)  # 1 px = 9525 EMU
prs.slide_height = Emu(1080 * 9525)

# PowerPoint default has 11 layouts; We'll use blank layout index 6 for custom building when needed.
blank_layout = prs.slide_layouts[6]
title_layout = prs.slide_layouts[0]
title_content_layout = prs.slide_layouts[1]

# ----------------------------------------------------------------------------------
# 1. TITLE SLIDE
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
# Background colour as Ground Forest (simulate scrim over image)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = GROUND_FOREST

# Add optional hero photo placeholder across slide (simulate user replace)
hero = add_placeholder_image_frame(slide, 0, 0, prs.slide_width, prs.slide_height, "Hero Photo / Field Ops Image")
hero.line.fill.background()

# scrim overlay (semi transparent)
scrim = add_scrim(slide, GROUND_FOREST, 0.7)

# Title
title_box = add_textbox(slide, MARGIN_L, Inches(2), Inches(9), Inches(2), "OWN THE GROUND.\nDELIVER THE FUTURE.", 
                        font_name=FONT_HEAD, font_size=Pt(72), bold=True, colour=VOLT_FIELD, align=PP_ALIGN.LEFT)

# Subtitle
subtitle_box = add_textbox(slide, MARGIN_L, Inches(4), Inches(9), Inches(1), 
                           "Ground Control × Nike Performance & Sustainability Theme", 
                           font_name=FONT_BODY, font_size=Pt(28), bold=False, colour=CHALK_WHITE, align=PP_ALIGN.LEFT)

add_logo_lockup(slide, left=MARGIN_L, top=prs.slide_height - Inches(0.6), height=Pt(24))

# ----------------------------------------------------------------------------------
# 2. AGENDA
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
# white background
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = CHALK_WHITE

# left vertical groundline swoosh strip
strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.5), prs.slide_height)
fill = strip.fill
fill.gradient()
grad = fill.gradient_stops
grad[0].color.rgb = GROUND_FOREST
grad[1].color.rgb = CANOPY_LIGHT
gs3 = grad.add_stop_position(0.9)
gs3.color.rgb = VOLT_FIELD
set_shape_line(strip, None, 0)

add_title(slide, "Agenda")

agenda_items = [
    "Welcome & Objectives",
    "Ground Control Impact Snapshot",
    "Performance Mindset in Field Ops",
    "Data & Automation Roadmap",
    "Action: Just Grow It Campaign"
]
y = MARGIN_T + Pt(70)
for i, item in enumerate(agenda_items, start=1):
    box = add_textbox(slide, MARGIN_L, y, Inches(10), Pt(36), f"{i}. {item}", 
                      font_name=FONT_BODY, font_size=Pt(28), bold=False, colour=CANOPY_LIGHT)
    y += Pt(42)

add_footer(slide)

# ----------------------------------------------------------------------------------
# 3. SECTION DIVIDER
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
# background Ground Forest
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = GROUND_FOREST

    # diagonal cut hero photo top right (approx by rotated rectangle)
hero = add_placeholder_image_frame(slide, prs.slide_width * 0.5, 0, prs.slide_width * 0.5, prs.slide_height * 0.6, "Section Image")
hero.rotation = 23
hero.left = int(prs.slide_width * 0.55)

# Section Title
add_textbox(slide, MARGIN_L, Inches(3), Inches(8), Inches(2), "SECTION TITLE", 
            font_name=FONT_HEAD, font_size=Pt(54), bold=True, colour=CHALK_WHITE, align=PP_ALIGN.LEFT)

# Volt underline
volt_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, MARGIN_L, Inches(4.8), Inches(3), Pt(6))
set_shape_fill(volt_line, VOLT_FIELD)
set_shape_line(volt_line, None, 0)

add_footer(slide)

# ----------------------------------------------------------------------------------
# 4. STANDARD CONTENT (TITLE + BODY)
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = CHALK_WHITE

add_title(slide, "Why Performance Thinking Matters in Field Ops")

body_text = (
    "Treat operational KPIs like athletic training targets.\n\n"
    "• Track progress: hectares restored, incidents prevented, uptime.\n"
    "• Celebrate PBs: record planting days, zero-incident streaks.\n"
    "• Coach & improve: feedback loops, wearable sensors, app alerts."
)
box = add_textbox(slide, MARGIN_L, MARGIN_T + Pt(80), Inches(10.5), Inches(4), body_text, 
                  font_name=FONT_BODY, font_size=Pt(20), bold=False, colour=CARBON_BLACK)
add_footer(slide)

# ----------------------------------------------------------------------------------
# 5. CONTENT + IMAGE RIGHT
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = CHALK_WHITE

add_title(slide, "Automation in the Field")

# text left
txt = (
    "Apps & sensors feed live progress data.\n"
    "Automated alerts when targets drift.\n"
    "Gamified tasks: earn Volt Badges for on-time completion."
)
box = add_textbox(slide, MARGIN_L, MARGIN_T + Pt(80), Inches(6), Inches(3), txt,
                  font_name=FONT_BODY, font_size=Pt(20), colour=CARBON_BLACK)

# image right
img_w = Inches(5)
img_h = Inches(3.5)
img_left = prs.slide_width - MARGIN_R - img_w
img_top = MARGIN_T + Pt(80)
add_placeholder_image_frame(slide, img_left, img_top, img_w, img_h, "Field App Screenshot")

add_footer(slide)

# ----------------------------------------------------------------------------------
# 6. DATA / CHART SLIDE (sample clustered column chart)
# ----------------------------------------------------------------------------------
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

slide = prs.slides.add_slide(blank_layout)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = CHALK_WHITE

add_title(slide, "Biodiversity Plots Restored vs Target")

chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('Actual', (80, 110, 150, 190))
chart_data.add_series('Target', (100, 120, 160, 200))

x, y, cx, cy = MARGIN_L, MARGIN_T + Pt(80), prs.slide_width - (MARGIN_L + MARGIN_R), Inches(4)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

# style colours manually (series order Actual=Canopy, Target=VOLT? Actually we want Target line)
# For columns we can't line; Instead we re-colour: Actual=CANOPY, Target=TRACKLINE
chart.series[0].format.fill.solid()
chart.series[0].format.fill.fore_color.rgb = CANOPY_LIGHT
chart.series[1].format.fill.solid()
chart.series[1].format.fill.fore_color.rgb = TRACKLINE_GREY

# Data labels optional off

add_footer(slide)

# ----------------------------------------------------------------------------------
# 7. KPI BIG NUMBER
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = GROUND_FOREST

# KPI number
add_textbox(slide, 0, Inches(3), prs.slide_width, Inches(1.5), "98%", 
            font_name=FONT_HEAD, font_size=Pt(120), bold=True, colour=VOLT_FIELD, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(4.5), prs.slide_width, Inches(1), "Biodiversity Target Delivered", 
            font_name=FONT_BODY, font_size=Pt(32), bold=False, colour=CHALK_WHITE, align=PP_ALIGN.CENTER)

add_footer(slide)

# ----------------------------------------------------------------------------------
# 8. QUOTE SLIDE
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = GROUND_FOREST

# large Volt watermark quote marks (approx by text)
add_textbox(slide, MARGIN_L, MARGIN_T, Inches(1), Inches(1), "“", 
            font_name=FONT_HEAD, font_size=Pt(120), bold=True, colour=VOLT_FIELD)

quote = (
    "Every metre we maintain is an opportunity to restore biodiversity.\n"
    "Every shift is a chance to beat our last performance."
)
add_textbox(slide, MARGIN_L + Inches(0.8), MARGIN_T + Pt(80), Inches(10), Inches(2.5), quote, 
            font_name=FONT_HEAD, font_size=Pt(40), bold=True, colour=CHALK_WHITE)

add_textbox(slide, MARGIN_L + Inches(0.8), MARGIN_T + Pt(220), Inches(10), Inches(1), "— Ground Control Team Lead", 
            font_name=FONT_SUB, font_size=Pt(24), bold=False, colour=CANOPY_LIGHT)

add_footer(slide)

# ----------------------------------------------------------------------------------
# 9. PHOTO COLLAGE
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = OVERCAST

# three angled frames stepping upward
w = Inches(3.5)
h = Inches(2)
gap = Inches(0.5)
left = MARGIN_L
top = Inches(3)
ph1 = add_placeholder_image_frame(slide, left, top, w, h, "Image 1")
ph1.rotation = -10
ph2 = add_placeholder_image_frame(slide, left + w + gap, top - Inches(0.4), w, h, "Image 2")
ph2.rotation = 5
ph3 = add_placeholder_image_frame(slide, left + (w + gap)*2, top - Inches(0.8), w, h, "Image 3")
ph3.rotation = 15

add_title(slide, "Field Highlights")

add_footer(slide)

# ----------------------------------------------------------------------------------
# 10. CLOSING / THANK YOU
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = CANOPY_LIGHT

# gradient tip (Volt triangle top right)
tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, prs.slide_width - Inches(4), 0, Inches(4), Inches(2.5))
set_shape_fill(tri, VOLT_FIELD)
tri.rotation = 180  # point towards slide
set_shape_line(tri, None, 0)

add_textbox(slide, MARGIN_L, Inches(3), Inches(8), Inches(1.5), "THANK YOU", 
            font_name=FONT_HEAD, font_size=Pt(72), bold=True, colour=CARBON_BLACK)

add_textbox(slide, MARGIN_L, Inches(4.5), Inches(8), Inches(1.5), "Questions? Contact: your.name@ground-control.co.uk", 
            font_name=FONT_BODY, font_size=Pt(28), bold=False, colour=CARBON_BLACK)

add_logo_lockup(slide, left=MARGIN_L, top=slide.height - Inches(0.6), height=Pt(24))

# ----------------------------------------------------------------------------------
# 11. COLOUR SWATCH REFERENCE (bonus)
# ----------------------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = CHALK_WHITE

add_title(slide, "Colour Swatches")

swatches = [
    ("Ground Forest", GROUND_FOREST),
    ("Canopy Light", CANOPY_LIGHT),
    ("Volt Field", VOLT_FIELD),
    ("Hi-Vis Amber", HIVIS_AMBER),
    ("Soil Brown", SOIL_BROWN),
    ("Carbon Black", CARBON_BLACK),
    ("Chalk White", CHALK_WHITE),
    ("Overcast", OVERCAST),
    ("Trackline Grey", TRACKLINE_GREY),
]

sw_w = Inches(1.5)
sw_h = Inches(1)
x = MARGIN_L
y = MARGIN_T + Pt(80)
cols = 5
for idx, (name, col) in enumerate(swatches):
    row = idx // cols
    col_idx = idx % cols
    sw = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + col_idx * (sw_w + Inches(0.25)),
                                y + row * (sw_h + Inches(0.6)), sw_w, sw_h)
    set_shape_fill(sw, col)
    set_shape_line(sw, CARBON_BLACK, 0.5)
    lbl = add_textbox(slide, sw.left, sw.top + sw_h + Pt(4), sw_w, Pt(20), name, 
                      font_name=FONT_BODY, font_size=Pt(12), colour=CARBON_BLACK, align=PP_ALIGN.CENTER)

add_footer(slide)


